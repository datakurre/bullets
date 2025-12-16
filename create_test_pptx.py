#!/usr/bin/env python3
"""
Create a simple test PPTX file for testing the import functionality.
Requires: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx not installed. Install with: pip install python-pptx")
    exit(1)

# Create a presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9 aspect ratio

# Slide 1: Title slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "Created for PPTX Import Testing"

# Slide 2: Title and content
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
content = slide2.placeholders[1]
title.text = "Bullet Points Example"
text_frame = content.text_frame
text_frame.text = "First bullet point"
p = text_frame.add_paragraph()
p.text = "Second bullet point"
p.level = 0
p = text_frame.add_paragraph()
p.text = "Third bullet point"
p.level = 0

# Slide 3: Title only
slide3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
txBox = slide3.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "Simple Text Slide"
tf.paragraphs[0].font.size = Pt(44)
tf.paragraphs[0].font.bold = True

# Slide 4: Multiple paragraphs
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide4.shapes.title
content = slide4.placeholders[1]
title.text = "Features"
text_frame = content.text_frame
text_frame.text = "Easy to use"
p = text_frame.add_paragraph()
p.text = "Markdown support"
p = text_frame.add_paragraph()
p.text = "PowerPoint import/export"

# Save the presentation
filename = 'test_presentation.pptx'
prs.save(filename)
print(f"Created {filename} successfully!")
print(f"You can now test importing this file in the Bullets application.")
